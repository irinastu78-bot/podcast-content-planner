"""
services/source_service.py
Извлечение чистого текста из разных источников:
- вставленный текст
- файлы: .txt, .md, .pdf, .docx, .pptx, .xlsx
- ссылка на веб-статью

Единый вход: extract_text(...). Любой источник приводится к строке,
которая затем подставляется в промпт как материал-источник.
"""

from __future__ import annotations
import io


# Ограничение, чтобы не отправить в модель гигантский текст.
MAX_CHARS = 12000


def _truncate(text: str) -> str:
    text = (text or "").strip()
    if len(text) > MAX_CHARS:
        return text[:MAX_CHARS] + "\n\n[...материал обрезан...]"
    return text


def from_plain_text(text: str) -> str:
    """Текст, вставленный пользователем напрямую."""
    return _truncate(text)


def from_txt(file_bytes: bytes) -> str:
    """.txt / .md — читаем как UTF-8 с запасным декодированием."""
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("cp1251", errors="ignore")
    return _truncate(text)


def from_pdf(file_bytes: bytes) -> str:
    """PDF — извлекаем текст постранично через pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    parts = [page.extract_text() or "" for page in reader.pages]
    return _truncate("\n".join(parts))


def from_docx(file_bytes: bytes) -> str:
    """Word .docx — параграфы и таблицы через python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return _truncate("\n".join(parts))


def from_pptx(file_bytes: bytes) -> str:
    """PowerPoint .pptx — текст из всех фигур слайдов через python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Слайд {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return _truncate("\n".join(parts))


def from_xlsx(file_bytes: bytes) -> str:
    """Excel .xlsx — содержимое ячеек по листам через openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"--- Лист: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return _truncate("\n".join(parts))


def from_url(url: str) -> str:
    """Веб-статья — извлекаем основной текст через trafilatura."""
    import trafilatura

    downloaded = trafilatura.fetch_url(url)
    if not downloaded:
        raise ValueError(f"Не удалось загрузить страницу: {url}")
    text = trafilatura.extract(
        downloaded,
        include_comments=False,
        include_tables=True,
        favor_recall=True,
    )
    if not text:
        raise ValueError(
            "Не удалось извлечь текст статьи. "
            "Попробуйте скопировать текст вручную."
        )
    return _truncate(text)


# Маппинг расширений файлов на обработчики.
_FILE_HANDLERS = {
    "txt": from_txt,
    "md": from_txt,
    "pdf": from_pdf,
    "docx": from_docx,
    "pptx": from_pptx,
    "xlsx": from_xlsx,
}

# Список расширений для виджета загрузки в Streamlit.
SUPPORTED_FILE_TYPES = list(_FILE_HANDLERS.keys())


def from_file(filename: str, file_bytes: bytes) -> str:
    """Определяет тип файла по расширению и вызывает нужный обработчик."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    handler = _FILE_HANDLERS.get(ext)
    if handler is None:
        raise ValueError(
            f"Формат .{ext} не поддерживается. "
            f"Доступны: {', '.join(SUPPORTED_FILE_TYPES)}"
        )
    return handler(file_bytes)


def extract_text(
    *,
    plain_text: str | None = None,
    filename: str | None = None,
    file_bytes: bytes | None = None,
    url: str | None = None,
) -> str:
    """
    Универсальная точка входа. Возвращает чистый текст из любого источника.
    Передавай ровно один источник.
    Если источника нет — вернётся пустая строка.
    """
    if plain_text and plain_text.strip():
        return from_plain_text(plain_text)
    if filename and file_bytes:
        return from_file(filename, file_bytes)
    if url and url.strip():
        return from_url(url.strip())
    return ""
